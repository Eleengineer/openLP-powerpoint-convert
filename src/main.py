#python3.6.8
targetFolder='\OneDrive\Documents\Power point song\\'       # where should we read from
import datetime
import multiprocessing
import os
import os.path
from xml.sax.saxutils import escape

import lxml.etree as etree

try:
  from pptx import Presentation
except ModuleNotFoundError:
  raise ModuleNotFoundError("Could not connect with powerpoint\nload with 'pip install python-pptx'")
home = os.path.expanduser('~')
#settings = {'layouts':{'content':0}}

def scanPresentation(targetPresentation):
  prs = Presentation(home + targetFolder + targetPresentation )
  root=etree.Element('root')
  song = etree.SubElement(root,'song',attrib={'version':'0.8'})
  properties = etree.SubElement(song,'properties')
  lyrics= etree.SubElement(song,'lyrics')
  titles = etree.SubElement(properties,'titles')
  title = etree.SubElement(titles ,'title')
  title.text = targetPresentation[:-5]
  verseCount= 0
  for slide in prs.slides:
    print('Checking Slide '+ str(prs.slides.index(slide)+1))
    if slide.slide_layout == prs.slide_layouts[0] and verseCount == 0:
      print('Found Title')
      #print(list(slide.shapes))
      if "text" in slide.shapes:
        print('slide has text shape')
        if 'Page' in slide.shapes.text:
          string=slide.shapes.title.text
          [int(s) for s in str.split() if s.isdigit()]
      save_path_file = slide.shapes.title.text
    else:# slide.slide_layout== prs.slide_layouts[6]:
      for shape in slide.shapes:
        if not shape.has_text_frame:
          continue
        #print(shape.text_frame.text)
        verseCount+=1
        verse= etree.SubElement(lyrics,'verse',attrib={'lang':'en','name':'v'+str(verseCount)})
        lines = etree.SubElement(verse,'lines')
        lines.text=''
        for paragraph in shape.text_frame.paragraphs:
          for run in paragraph.runs:
            lines.text= str(lines.text) + escape(run.text, {"’" : "&apos","‘" : "&apos",'♂':'\n'}) + "\n"
        #else:
      #print('no match')
   
  #'http://openlyrics.info/namespace/2009/song') 
  xml.setAttribute('modifiedDate',datetime.datetime.now().strftime())  
  song.appendChild(xml)
  titles = song.createElement('titles')
  titles.appendChild(title)
  
  properties.setAttribute('xmlns', 'http://openlyrics.info/namespace/2009/song')
  
  xml_str = etree.tounicode(song)
  del root, prs
  with open('xml\\'+ targetPresentation[:-5] + '.xml', "w") as f:
    f.write(xml_str)
  
if __name__ == '__main__':
  for file in os.listdir(home + targetFolder):
    print(file)
    if not file.endswith('.pptx'):
      continue
    #print('good')
    scanPresentation(file)
